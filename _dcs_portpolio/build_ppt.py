#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""포트폴리오 PPTX 생성 스크립트. /working/result 의 md 문서를 기반으로 슬라이드 구성."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

FONT = "맑은 고딕"
NAVY = RGBColor(0x1F, 0x29, 0x37)
BLUE = RGBColor(0x2F, 0x6F, 0xED)
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x55, 0x55, 0x55)
GREEN = RGBColor(0x1E, 0x7A, 0x34)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def set_text(tf, text, size=18, color=NAVY, bold=False, align=PP_ALIGN.LEFT, font=FONT):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return p


def add_textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def add_footer(slide, page_no, total):
    tf = add_textbox(slide, Inches(0.5), Inches(7.08), Inches(6), Inches(0.35))
    set_text(tf, "컨테이너 오케스트레이션 & 통합 관제 플랫폼", size=10, color=RGBColor(0x99, 0x99, 0x99))
    tf2 = add_textbox(slide, Inches(12.3), Inches(7.08), Inches(0.7), Inches(0.35))
    set_text(tf2, f"{page_no}", size=10, color=RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.RIGHT)


def title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    bar.shadow.inherit = False
    tf = add_textbox(slide, Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.6))
    set_text(tf, title, size=26, color=WHITE, bold=True)
    if subtitle:
        tf2 = add_textbox(slide, Inches(0.5), Inches(0.68), Inches(12.3), Inches(0.4))
        set_text(tf2, subtitle, size=13, color=RGBColor(0xCB, 0xD5, 0xE1))


def bullets(slide, left, top, width, height, items, size=16, color=NAVY, line_gap=6):
    tf = add_textbox(slide, left, top, width, height)
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(line_gap)
        p.level = level
        r = p.add_run()
        prefix = "•  " if level == 0 else "-  "
        r.text = prefix + text
        r.font.size = Pt(size - level * 2)
        r.font.color.rgb = color
        r.font.name = FONT
    return tf


def box(slide, left, top, width, height, text, fill=BLUE, text_color=WHITE, size=13, bold=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = slide.shapes.add_shape(shape, left, top, width, height)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = fill
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = text_color
    r.font.name = FONT
    return sp


def arrow(slide, x1, y1, x2, y2, color=GRAY, width=1.5):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    conn.line.end_arrowhead = True if hasattr(conn.line, "end_arrowhead") else None
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    return conn


# ---------------------------------------------------------------- Slide 1: Title
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
tf = add_textbox(s, Inches(1), Inches(2.6), Inches(11.3), Inches(1.3))
set_text(tf, "컨테이너 오케스트레이션 &", size=40, color=WHITE, bold=True)
tf2 = add_textbox(s, Inches(1), Inches(3.25), Inches(11.3), Inches(1.0))
set_text(tf2, "통합 관제 플랫폼", size=40, color=WHITE, bold=True)
tf3 = add_textbox(s, Inches(1), Inches(4.35), Inches(11.3), Inches(0.8))
set_text(tf3, "분산 컨테이너 인스턴스 관제 · 파일 릴레이 · 모니터링 · 알림 통합 — 설계부터 운영까지",
         size=17, color=RGBColor(0xCB, 0xD5, 0xE1))
tf4 = add_textbox(s, Inches(1), Inches(6.6), Inches(11.3), Inches(0.5))
set_text(tf4, "Backend / Infra Portfolio", size=13, color=RGBColor(0x93, 0xA3, 0xB8))

# ---------------------------------------------------------------- Slide 2: 배경/문제정의
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
title_bar(s, "프로젝트 배경", "여러 서버에 흩어진 컨테이너를 통합 관리해야 하는 문제")
bullets(s, Inches(0.7), Inches(1.6), Inches(11.9), Inches(5),
        [
            "업무 컨테이너 인스턴스가 여러 물리 서버에 분산 배치되어, 서버별 SSH 접속 후 " +
            "docker-compose up/down을 수동 실행 — 상태를 한눈에 파악할 방법이 없었음",
            "인스턴스 간 파일 연동(발신/수신)이 필요했지만, 검증되지 않은 임시 스크립트에 의존",
            "인프라 모니터링(CPU/메모리/컨테이너 상태)은 있었지만, 업무 지표(처리 건수 등)는 " +
            "별도로 확인할 방법이 없었음",
            "장애(컨테이너 다운) 발생 시 담당자가 화면을 계속 보고 있어야만 인지 가능",
        ], size=18, line_gap=18)
add_footer(s, 2, 13)

# ---------------------------------------------------------------- Slide 3: 아키텍처 다이어그램
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
title_bar(s, "시스템 아키텍처", "Manager → Server → Agent 3계층 + 파일 릴레이 + 모니터링/알림")

# Row 1: Manager
box(s, Inches(0.6), Inches(1.5), Inches(3.0), Inches(0.7), "관제 웹 콘솔\n(Manager)", fill=BLUE)
# Row 2: Server
box(s, Inches(0.6), Inches(2.55), Inches(3.0), Inches(0.7), "중계 서버\n(Server)", fill=BLUE)
# Row 3: Agents x2
box(s, Inches(0.2), Inches(3.6), Inches(1.9), Inches(0.65), "Host Agent #1", fill=NAVY, size=11)
box(s, Inches(2.3), Inches(3.6), Inches(1.9), Inches(0.65), "Host Agent #2", fill=NAVY, size=11)
# Row 4: containers
box(s, Inches(0.1), Inches(4.5), Inches(1.0), Inches(0.55), "업무\n컨테이너", fill=GRAY, size=9)
box(s, Inches(1.2), Inches(4.5), Inches(1.0), Inches(0.55), "업무\n컨테이너", fill=GRAY, size=9)
box(s, Inches(2.4), Inches(4.5), Inches(1.0), Inches(0.55), "업무\n컨테이너", fill=GRAY, size=9)
box(s, Inches(3.5), Inches(4.5), Inches(1.0), Inches(0.55), "업무\n컨테이너", fill=GRAY, size=9)

# File relay
box(s, Inches(5.0), Inches(3.6), Inches(2.6), Inches(0.65), "File Relay Sender\n(각 인스턴스 sidecar)", fill=RGBColor(0x7C, 0x3A, 0xED), size=11)
box(s, Inches(5.0), Inches(4.5), Inches(2.6), Inches(0.75), "File Relay Receiver\n(자체 TCP 프로토콜)", fill=RGBColor(0x7C, 0x3A, 0xED), size=11)

# Monitoring
box(s, Inches(8.2), Inches(1.5), Inches(2.3), Inches(0.6), "cAdvisor / node_exporter", fill=RGBColor(0x0D, 0x94, 0x88), size=10)
box(s, Inches(8.2), Inches(2.35), Inches(2.3), Inches(0.6), "Prometheus", fill=RGBColor(0x0D, 0x94, 0x88), size=11)
box(s, Inches(8.2), Inches(3.2), Inches(2.3), Inches(0.6), "Grafana 대시보드", fill=RGBColor(0x0D, 0x94, 0x88), size=11)

# Notification
box(s, Inches(11.0), Inches(1.5), Inches(1.9), Inches(0.85), "개인 메신저\n알림 API", fill=RGBColor(0xD6, 0x45, 0x45), size=11)

# Arrows
arrow(s, Inches(2.1), Inches(2.2), Inches(2.1), Inches(2.55))
arrow(s, Inches(1.6), Inches(3.25), Inches(1.6), Inches(3.6))
arrow(s, Inches(3.3), Inches(3.25), Inches(3.3), Inches(3.6))
arrow(s, Inches(0.7), Inches(4.25), Inches(0.7), Inches(4.5))
arrow(s, Inches(1.8), Inches(4.25), Inches(1.8), Inches(4.5))
arrow(s, Inches(6.3), Inches(4.25), Inches(6.3), Inches(4.5))
arrow(s, Inches(3.6), Inches(1.85), Inches(5.0), Inches(4.8))
arrow(s, Inches(3.6), Inches(2.9), Inches(9.3), Inches(2.65))
arrow(s, Inches(2.85), Inches(2.55), Inches(11.9), Inches(1.75))

note = add_textbox(s, Inches(0.7), Inches(5.5), Inches(11.9), Inches(1.6))
bullets(s, Inches(0.7), Inches(5.5), Inches(11.9), Inches(1.6),
        [
            "Manager는 어떤 물리서버에 어떤 인스턴스가 있는지 몰라도 됨 — Server가 라우팅을 담당",
            "File Relay는 HTTP가 아닌 자체 TCP 바이너리 프로토콜로 구현 (대용량 스트리밍 + pull 방식 수신 지원)",
            "Server가 상태를 주기적으로 폴링하며 up/down 전이를 감지해 알림 API 호출",
        ], size=13, line_gap=6)
add_footer(s, 3, 13)

# ---------------------------------------------------------------- Slide 4: 핵심 성과
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
title_bar(s, "핵심 성과")
achievements = [
    "가동/중지, 신규 등록, 삭제를 웹 UI에서 원클릭으로 처리 — 서버별 SSH 접속 불필요",
    "인스턴스 확장 가능한 구조 설계 (Server가 DB 기반 라우팅 테이블로 물리 서버 위치 관리)",
    "자체 TCP 프로토콜로 파일 릴레이 시스템을 처음부터 구현, 운영 실데이터 보호 로직 적용",
    "Prometheus/Grafana 인프라 모니터링에 업무 지표를 자연스럽게 통합",
    "상태 폴링 로직에 전이 감지를 추가해 장애 발생 시 개인 메신저로 즉시 알림",
    "동기 블로킹 API를 비동기로 리팩터링해 UI 응답시간을 수십 초 → 0.2초로 단축",
]
top = Inches(1.7)
for i, a in enumerate(achievements):
    row = i
    y = top + Inches(0.85) * row
    num = box(s, Inches(0.7), y, Inches(0.6), Inches(0.6), str(i + 1), fill=BLUE, size=18)
    tf = add_textbox(s, Inches(1.5), y + Inches(0.03), Inches(11.0), Inches(0.75))
    set_text(tf, a, size=16, color=NAVY)
add_footer(s, 4, 13)


# ---------------------------------------------------------------- Case study slides
cases = [
    {
        "title": "Case 1 · JS 동적 폼에서만 발생하는 403 CSRF 오류",
        "sit": "관리 화면 '실행' 버튼만 항상 403 발생 — 서버 렌더링 폼(삭제 버튼 등)은 정상",
        "cause": "document.createElement('form')으로 즉석 생성한 폼에는 Spring Security가 " +
                 "자동으로 넣어주는 CSRF 히든 필드가 없었음",
        "fix": "페이지에 이미 렌더링된 다른 폼에서 CSRF 토큰 값을 그대로 읽어와 동적 폼에 추가",
        "lesson": "재현 안 될 때 테스트 도구 문제로 성급히 결론짓지 않고, 실사용자 환경 증거(접속 로그)로 재검증",
    },
    {
        "title": "Case 2 · 컨테이너 기본 시간대(UTC)로 인한 9시간 오차",
        "sit": "생성 파일의 타임스탬프가 실제 생성 시각과 정확히 9시간(UTC-KST) 차이",
        "cause": "컨테이너 베이스 이미지에 TZ 미설정 → LocalDateTime.now()가 UTC 기준으로 동작",
        "fix": "ZoneId.of(\"Asia/Seoul\")을 코드에 명시적으로 고정, 컨테이너 설정에 의존하지 않도록 변경",
        "lesson": "환경설정 의존 대신 비즈니스 로직에 필요한 시간대는 코드 레벨에서 명시적으로 고정",
    },
    {
        "title": "Case 3 · Prometheus 집계 함수 오류로 상태값 오표시",
        "sit": "정상 실행 중인 컨테이너 상태가 '1(UP)'이 아니라 '2'로 표시",
        "cause": "재시작 시 cAdvisor가 신/구 라벨 조합을 일시적으로 동시 노출 → count()가 라벨 " +
                 "조합 개수를 셈",
        "fix": "존재 여부 확인 용도이므로 count 대신 group으로 변경 (라벨 카디널리티 변화에 안전)",
        "lesson": "'몇 개 있는지'가 아니라 '있는지 없는지'를 물을 땐 count보다 group이 안전",
    },
    {
        "title": "Case 4 · 인프라 모니터링에 업무 지표 통합",
        "sit": "CPU/메모리/상태만 보여주던 대시보드에 '오늘 처리 건수'도 같이 보고 싶다는 요구",
        "cause": "Prometheus/Grafana 표준 스택엔 업무 지표를 낼 방법이 기본으로 없음",
        "fix": "node_exporter textfile collector + cron 스크립트로 커스텀 지표 노출, PromQL " +
               "'and on(label)'로 유령 데이터(미존재 인스턴스) 필터링",
        "lesson": "'별도 시스템 필요'라 단정 말고 기존 스택의 확장 포인트를 먼저 탐색",
    },
    {
        "title": "Case 5 · 재기동에도 살아남는 OAuth 토큰 생명주기 관리",
        "sit": "장애 알림용 메신저 API는 6시간 액세스 토큰 + 60일 리프레시 토큰(회전 가능) 구조",
        "cause": "단순 메모리 보관 시 재기동마다 재인증 필요, 토큰 회전을 놓치면 알림 영구 중단",
        "fix": "액세스 토큰은 5시간 주기 자동 갱신, 리프레시 토큰은 회전 시마다 파일에 원자적으로 재저장",
        "lesson": "외부 API 토큰 생명주기는 '재기동돼도 동작하는가'까지 설계에 반영",
    },
    {
        "title": "Case 6 · 동기 블로킹 API를 비동기로 리팩터링",
        "sit": "파일 N개를 M초 간격 생성하는 기능에서 (N-1)×M초 동안 브라우저가 그대로 멈춤",
        "cause": "요청 스레드가 생성 루프의 Thread.sleep까지 전부 블로킹한 뒤 응답 — HTTP 체인 " +
                 "전체가 동기로 연결",
        "fix": "유효성 검증만 동기 처리 후 즉시 '접수됨' 응답, 실제 작업은 백그라운드 스레드풀에 위임",
        "lesson": "응답시간 수십 초 → 0.2초. '접수됨'과 '완료됨'을 화면에서도 명확히 구분",
    },
    {
        "title": "Case 7 · 자체 TCP 프로토콜 기반 파일 릴레이 (제로베이스 설계)",
        "sit": "인스턴스↔중앙서버 양방향 파일 릴레이, 대용량 스트리밍 + pull 방식 수신 필요",
        "cause": "HTTP로도 가능하나 가벼운 자체 프로토콜이 더 적합하다고 판단",
        "fix": "1바이트 opType + writeUTF/writeLong 프레이밍, .part 임시파일 후 원자적 rename, " +
               "서비스 최초 기동 시 기존 파일을 베이스라인으로 스냅샷해 영구 보호",
        "lesson": "운영 중인 실데이터에 새 자동화를 얹을 땐 최소 영향 원칙을 항상 먼저 고려",
    },
]

for idx, c in enumerate(cases):
    s = prs.slides.add_slide(BLANK)
    add_bg(s, WHITE)
    title_bar(s, c["title"])
    labels = [("상황", c["sit"], RGBColor(0x4A, 0x7F, 0xD6)),
              ("원인", c["cause"], RGBColor(0xD6, 0x45, 0x45)),
              ("해결", c["fix"], GREEN),
              ("배운 점", c["lesson"], RGBColor(0x7C, 0x3A, 0xED))]
    y = Inches(1.6)
    for label, text, color in labels:
        tag = box(s, Inches(0.7), y, Inches(1.3), Inches(0.55), label, fill=color, size=13)
        tf = add_textbox(s, Inches(2.2), y - Inches(0.02), Inches(10.3), Inches(1.0))
        set_text(tf, text, size=15, color=NAVY)
        y += Inches(1.15)
    add_footer(s, 5 + idx, 13)

# ---------------------------------------------------------------- Tech stack slide
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)
title_bar(s, "기술 스택")
cols = [
    ("Backend", ["Java 8 / 21, Spring Boot 2.7", "Spring MVC / Security / Data JPA",
                 "Thymeleaf", "REST API (API Key 인증)",
                 "Raw TCP 소켓 프로그래밍", "ExecutorService / @Scheduled 비동기"]),
    ("인프라 · 운영", ["Docker / Docker Compose", "사설 컨테이너 레지스트리",
                    "Linux(CentOS), systemd, cron", "LVM/XFS, firewalld",
                    "MariaDB"]),
    ("모니터링 · 연동", ["Prometheus (PromQL)", "Grafana (대시보드/API)",
                    "cAdvisor, node_exporter", "OAuth2 스타일 토큰 연동",
                    "외부 메신저 알림 API"]),
]
col_w = Inches(3.9)
for i, (head, items) in enumerate(cols):
    x = Inches(0.6) + col_w * i + Inches(0.15) * i
    box(s, x, Inches(1.6), col_w, Inches(0.55), head, fill=NAVY, size=15)
    bullets(s, x, Inches(2.3), col_w, Inches(4.0), items, size=13, line_gap=10)
add_footer(s, 12, 13)

# ---------------------------------------------------------------- Closing slide
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
tf = add_textbox(s, Inches(1), Inches(3.0), Inches(11.3), Inches(1.0))
set_text(tf, "감사합니다", size=36, color=WHITE, bold=True)
tf2 = add_textbox(s, Inches(1), Inches(3.9), Inches(11.3), Inches(0.8))
set_text(tf2, "설계부터 운영/트러블슈팅까지 전 과정을 직접 수행한 경험을 담았습니다.",
         size=15, color=RGBColor(0xCB, 0xD5, 0xE1))
add_footer(s, 13, 13)

prs.save("/working/result/portfolio.pptx")
print("saved")
